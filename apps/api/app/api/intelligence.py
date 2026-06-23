"""
Layer 1 Intelligence Report API
POST /intelligence/generate  — run connectors for an entity, build graph, generate cited dossier
GET  /intelligence/{report_id} — retrieve a generated intelligence report
GET  /intelligence/            — list recent intelligence reports
"""
from fastapi import APIRouter, Depends, BackgroundTasks, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from sqlalchemy import text
from typing import Optional, Dict, Any, List
import io
from app.db.session import get_db
from app.services.intelligence_service import generate_intelligence_report, get_intelligence_report, list_intelligence_reports

try:
    from app.connectors.browser_research_agent import research_entity_browser, detect_jurisdiction
    _BROWSER_AVAILABLE = True
except ImportError:
    _BROWSER_AVAILABLE = False

try:
    from app.connectors.apollo_connector import (
        search_people, enrich_organization, fetch_org_chart, search_organization,
    )
    _APOLLO_AVAILABLE = True
except ImportError:
    _APOLLO_AVAILABLE = False

try:
    from app.connectors.private_company_connector import (
        fetch_private_company_intel,
        search_opencorporates,
        search_gleif,
        search_fincen_entities,
    )
    _PRIV_CO_AVAILABLE = True
except ImportError:
    _PRIV_CO_AVAILABLE = False

try:
    from app.services.pdf_service import generate_report_pdf
    _PDF_AVAILABLE = True
except ImportError:
    _PDF_AVAILABLE = False

router = APIRouter(prefix="/intelligence")


@router.post("/generate")
def generate_report(
    entity_name: str,
    entity_type: str = "org",
    ticker: Optional[str] = None,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
):
    """
    Kick off a Layer 1 Entity Network Intelligence Report.
    Runs connectors (SEC, FEC, FARA, USASpending, LDA, OFAC, CourtListener),
    writes relationships + evidence, assembles report JSON, generates GPT narrative.
    Returns report_id immediately; report builds synchronously for demo.
    """
    report = generate_intelligence_report(db, entity_name=entity_name, entity_type=entity_type, ticker=ticker)
    return report


@router.post("/browser-research")
def browser_research(
    entity_name: str,
    jurisdiction: Optional[str] = None,
    context: str = "",
    deep_dive: bool = False,
):
    """
    Run browser-based research on an entity using public registries, news, and
    government sources for any jurisdiction (used as fallback for non-US entities
    or for deep dives on holding companies, vehicles, financials).

    Examples:
      POST /intelligence/browser-research?entity_name=Aeropuertos+Argentina+2000&jurisdiction=argentina
      POST /intelligence/browser-research?entity_name=Mercado+Libre&deep_dive=true
    """
    if not _BROWSER_AVAILABLE:
        raise HTTPException(503, "Browser research agent not available")
    detected_jur = jurisdiction or detect_jurisdiction(entity_name, context)
    result = research_entity_browser(
        entity_name,
        jurisdiction=detected_jur,
        context=context,
        max_sources=6 if deep_dive else 4,
        deep_dive=deep_dive,
    )
    return result


@router.get("/")
def list_reports(limit: int = 20, db: Session = Depends(get_db)):
    return list_intelligence_reports(db, limit=limit)


@router.get("/apollo/org")
def apollo_org(name: str, domain: str = ""):
    """
    Enrich an organization using Apollo.io.
    Returns headcount, revenue range, technologies, founding year, etc.

    Strategy: name-based search first (uses Apollo's internal domain resolution),
    then direct domain enrichment if domain is explicitly provided.
    """
    if not _APOLLO_AVAILABLE:
        raise HTTPException(503, "Apollo connector not available")

    result = {}
    # Explicit domain takes priority
    if domain:
        result = enrich_organization(domain=domain)
    # Name-based search is most accurate (Apollo resolves to the correct entity)
    if not result and name:
        result = search_organization(name)
    # Last resort: guess domain
    if not result and name:
        slug = name.lower().replace(" ", "").replace(".", "").replace(",", "")
        result = enrich_organization(domain=f"{slug}.com")

    return result or {"message": "No Apollo data found", "hint": "Check APOLLO_API_KEY and try adding domain= parameter"}


@router.get("/apollo/people")
def apollo_people(organization: str, limit: int = 20):
    """
    Fetch key executives and employees of an organization via Apollo.io.
    Returns up to `limit` people with name, title, email, LinkedIn URL.
    """
    if not _APOLLO_AVAILABLE:
        raise HTTPException(503, "Apollo connector not available")
    return search_people(organization=organization, limit=limit)


@router.get("/apollo/orgchart")
def apollo_orgchart(organization: str):
    """
    Fetch the C-suite + VP-level org chart for an organization via Apollo.io.
    """
    if not _APOLLO_AVAILABLE:
        raise HTTPException(503, "Apollo connector not available")
    return fetch_org_chart(organization)


@router.post("/apollo/enrich")
def apollo_enrich(
    entity_name: str,
    domain: str = "",
    include_people: bool = True,
):
    """
    Full Apollo enrichment: org profile + executives + key people search.

    Free plan: org enrichment only (by domain preferred).
    Paid plan: + people/org chart search.
    """
    if not _APOLLO_AVAILABLE:
        raise HTTPException(503, "Apollo connector not available — set APOLLO_API_KEY in .env")

    # Domain lookup first (most accurate on free tier)
    org_data = {}
    if domain:
        org_data = enrich_organization(domain=domain)
    if not org_data:
        # Try to derive domain from name
        slug = entity_name.lower().replace(" ", "").replace(".", "").replace(",", "")
        org_data = enrich_organization(domain=f"{slug}.com")
    if not org_data:
        org_data = search_organization(entity_name)

    people_data = []
    people_plan_note = ""
    if include_people:
        people_data = fetch_org_chart(entity_name)
        if not people_data:
            people_plan_note = "People search requires Apollo paid plan (free tier: org enrichment only)"

    return {
        "entity_name":       entity_name,
        "organization":      org_data,
        "key_people":        people_data,
        "total_people":      len(people_data),
        "plan_note":         people_plan_note or None,
        "data_coverage":     {
            "org_enrichment":  bool(org_data),
            "people_search":   bool(people_data),
            "free_tier_only":  not bool(people_data),
        },
    }


@router.get("/private-co/search")
def private_co_search(name: str, jurisdiction: str = ""):
    """
    Full private company enrichment: OpenCorporates + GLEIF + FinCEN + FDIC.
    """
    if not _PRIV_CO_AVAILABLE:
        raise HTTPException(503, "Private company connector not available")
    return fetch_private_company_intel(name, jurisdiction=jurisdiction)


@router.get("/private-co/opencorporates")
def oc_search(name: str, jurisdiction: str = "", limit: int = 5):
    """Search OpenCorporates for company registrations globally."""
    if not _PRIV_CO_AVAILABLE:
        raise HTTPException(503, "Private company connector not available")
    return search_opencorporates(name, jurisdiction=jurisdiction, limit=limit)


@router.get("/private-co/gleif")
def gleif_search(name: str, limit: int = 3):
    """Search GLEIF for Legal Entity Identifiers."""
    if not _PRIV_CO_AVAILABLE:
        raise HTTPException(503, "Private company connector not available")
    return search_gleif(name, limit=limit)


@router.get("/private-co/fincen")
def fincen_search(name: str):
    """Search FinCEN financial institution registry."""
    if not _PRIV_CO_AVAILABLE:
        raise HTTPException(503, "Private company connector not available")
    return search_fincen_entities(name)


@router.get("/{report_id}")
def get_report(report_id: int, db: Session = Depends(get_db)):
    report = get_intelligence_report(db, report_id)
    if not report:
        raise HTTPException(404, "Intelligence report not found")
    return report


@router.get("/{report_id}/pdf")
def download_report_pdf(report_id: int, db: Session = Depends(get_db)):
    """
    Download an intelligence report as a polished PDF.
    Returns a binary PDF file suitable for direct browser download.
    """
    if not _PDF_AVAILABLE:
        raise HTTPException(503, "PDF export unavailable — install reportlab")

    report = get_intelligence_report(db, report_id)
    if not report:
        raise HTTPException(404, "Intelligence report not found")

    try:
        pdf_bytes = generate_report_pdf(report)
    except Exception as exc:
        raise HTTPException(500, f"PDF generation failed: {exc}")

    entity_slug = (report.get("entity_name") or "report").lower().replace(" ", "_")[:40]
    filename    = f"intel_{entity_slug}_{report_id}.pdf"

    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/{report_id}/word")
def download_report_word(report_id: int, db: Session = Depends(get_db)):
    """Download intelligence report as Word (.docx)."""
    try:
        from docx import Document as DocxDocument
        from docx.shared import Pt, RGBColor
    except ImportError:
        raise HTTPException(503, "python-docx not installed")

    report = get_intelligence_report(db, report_id)
    if not report:
        raise HTTPException(404, "Intelligence report not found")

    doc = DocxDocument()
    doc.add_heading(f"Intelligence Report: {report.get('entity_name', 'Unknown')}", 0)
    doc.add_paragraph(f"Generated: {report.get('created_at', '')}  |  Type: {report.get('entity_type', '')}")
    doc.add_paragraph("")

    for section in report.get("sections", []):
        doc.add_heading(section.get("title", "Section"), level=1)
        doc.add_paragraph(section.get("summary", ""))
        for claim in section.get("claims", []):
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(f"[{claim.get('confidence', 'ANALYTICAL')}] ").bold = True
            p.add_run(claim.get("text", ""))
            if claim.get("source_url"):
                p.add_run(f"  ({claim['source_url']})")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    entity_slug = (report.get("entity_name") or "report").lower().replace(" ", "_")[:40]
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="intel_{entity_slug}_{report_id}.docx"'},
    )


@router.get("/{report_id}/excel")
def download_report_excel(report_id: int, db: Session = Depends(get_db)):
    """Download intelligence report as Excel (.xlsx) — claims as rows."""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        raise HTTPException(503, "openpyxl not installed")

    report = get_intelligence_report(db, report_id)
    if not report:
        raise HTTPException(404, "Intelligence report not found")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Intelligence Report"

    # Header row
    headers = ["Section", "Claim", "Confidence", "Source URL"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="4F46E5")
        cell.alignment = Alignment(wrap_text=True)

    row = 2
    for section in report.get("sections", []):
        title = section.get("title", "")
        for claim in section.get("claims", []):
            ws.cell(row=row, column=1, value=title)
            ws.cell(row=row, column=2, value=claim.get("text", ""))
            ws.cell(row=row, column=3, value=claim.get("confidence", "ANALYTICAL"))
            ws.cell(row=row, column=4, value=claim.get("source_url", ""))
            row += 1

    ws.column_dimensions["A"].width = 30
    ws.column_dimensions["B"].width = 80
    ws.column_dimensions["C"].width = 15
    ws.column_dimensions["D"].width = 50

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    entity_slug = (report.get("entity_name") or "report").lower().replace(" ", "_")[:40]
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="intel_{entity_slug}_{report_id}.xlsx"'},
    )


@router.get("/{report_id}/powerpoint")
def download_report_pptx(report_id: int, db: Session = Depends(get_db)):
    """Download intelligence report as PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor as PptxRGB
    except ImportError:
        raise HTTPException(503, "python-pptx not installed")

    report = get_intelligence_report(db, report_id)
    if not report:
        raise HTTPException(404, "Intelligence report not found")

    prs = Presentation()
    blank_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Intelligence Report"
    slide.placeholders[1].text = f"{report.get('entity_name', 'Unknown')} — {report.get('created_at', '')[:10]}"

    # One slide per section (max 6 bullet points per slide)
    for section in report.get("sections", []):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = section.get("title", "Section")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        claims = section.get("claims", [])[:6]
        for i, claim in enumerate(claims):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"[{claim.get('confidence','?')}] {claim.get('text','')[:200]}"
            p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    entity_slug = (report.get("entity_name") or "report").lower().replace(" ", "_")[:40]
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="intel_{entity_slug}_{report_id}.pptx"'},
    )
